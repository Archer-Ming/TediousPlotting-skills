"""
Conversion utilities for the frame-skill. Run these AFTER you have an SVG ready.

Three public functions:
    svg_to_png(svg_path, png_path, dpi=400)
    svg_to_pptx(svg_path, pptx_path, slide_w_in=13.333, slide_h_in=7.5)
    validate_svg(svg_path) -> list[str]   # returns list of self-check problems

Dependencies:
    pip install cairosvg python-pptx Pillow

CLI usage:
    python render_helpers.py validate  reference/style_01_layered_network.svg
    python render_helpers.py to_png    reference/style_01_layered_network.svg out.png --dpi 400
    python render_helpers.py to_pptx   reference/style_01_layered_network.svg out.pptx
    python render_helpers.py to_all    reference/style_01_layered_network.svg     # writes out.png and out.pptx next to it
"""
from __future__ import annotations
import argparse
import os
import re
import sys
import tempfile
from pathlib import Path
from xml.etree import ElementTree as ET


#SVG → PNG
def svg_to_png(svg_path: str | Path, png_path: str | Path, dpi: int = 400) -> Path:
    """Render an SVG to PNG at the given DPI.

    DPI here means "output pixels per logical inch of the SVG viewBox at 96-dpi base".
    For a viewBox 1600x900 and dpi=400, output is roughly (1600 * 400/96) px wide.
    """
    import cairosvg  # imported lazily so the validator works without it
    svg_path = Path(svg_path)
    png_path = Path(png_path)
    png_path.parent.mkdir(parents=True, exist_ok=True)

    # Compute target pixel width from viewBox + dpi
    width_px = _target_width_from_viewbox(svg_path, dpi)
    cairosvg.svg2png(
        url=str(svg_path),
        write_to=str(png_path),
        output_width=width_px,
    )
    return png_path


def _target_width_from_viewbox(svg_path: Path, dpi: int) -> int:
    """Compute the pixel width that produces the requested DPI for an SVG.

    SVG units are nominally 96 user-units per inch. We scale so 1 user-unit
    becomes (dpi / 96) device pixels.
    """
    tree = ET.parse(svg_path)
    root = tree.getroot()
    vb = root.attrib.get("viewBox")
    if vb:
        parts = vb.replace(",", " ").split()
        if len(parts) == 4:
            _, _, w, _ = parts
            return max(64, int(float(w) * dpi / 96))
    # fallback: use width attribute if present
    w_attr = root.attrib.get("width", "1600")
    w_attr = re.sub(r"[^\d.]", "", w_attr) or "1600"
    return max(64, int(float(w_attr) * dpi / 96))


#SVG → PPTX
# Strategy: we don't try to parse SVG into native PPTX shapes (that fails
# for any non-trivial diagram and breaks every time the SVG is edited).
# Instead, we render the SVG to a high-resolution PNG and embed that PNG
# centered on a single slide. This guarantees the diagram in PowerPoint
# looks IDENTICAL to the SVG, and editing is done in the SVG source.

def svg_to_pptx(
    svg_path: str | Path,
    pptx_path: str | Path,
    slide_w_in: float = 13.333,   # 16:9 widescreen default
    slide_h_in: float = 7.5,
    dpi: int = 400,
    title: str | None = None,
) -> Path:
    """Render the SVG to a high-DPI PNG, then embed it centered on one slide."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from PIL import Image

    svg_path = Path(svg_path)
    pptx_path = Path(pptx_path)
    pptx_path.parent.mkdir(parents=True, exist_ok=True)

    # 1. Render PNG to a temp file
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
        tmp_png = Path(tmp.name)
    try:
        svg_to_png(svg_path, tmp_png, dpi=dpi)

        # 2. Read PNG size to compute correct aspect ratio
        with Image.open(tmp_png) as im:
            img_w_px, img_h_px = im.size
        img_aspect = img_w_px / img_h_px

        # 3. Set up a blank presentation
        prs = Presentation()
        prs.slide_width  = Inches(slide_w_in)
        prs.slide_height = Inches(slide_h_in)
        blank = prs.slide_layouts[6]  # fully blank layout
        slide = prs.slides.add_slide(blank)

        # 4. Compute the largest centered rectangle on the slide that
        #    matches img_aspect, with a small margin.
        margin_in = 0.4
        avail_w_in = slide_w_in - 2 * margin_in
        avail_h_in = slide_h_in - 2 * margin_in
        # If a title is requested, reserve top space
        if title:
            avail_h_in -= 0.6
        slide_aspect = avail_w_in / avail_h_in
        if img_aspect >= slide_aspect:
            # image is wider relative to slide → fit width
            pic_w_in = avail_w_in
            pic_h_in = avail_w_in / img_aspect
        else:
            pic_h_in = avail_h_in
            pic_w_in = avail_h_in * img_aspect
        left_in = (slide_w_in - pic_w_in) / 2
        top_offset = (0.6 if title else 0) + margin_in
        top_in = top_offset + (avail_h_in - pic_h_in) / 2

        # 5. Optional title above the image
        if title:
            from pptx.util import Emu
            tx = slide.shapes.add_textbox(Inches(margin_in), Inches(margin_in),
                                          Inches(slide_w_in - 2 * margin_in), Inches(0.5))
            tf = tx.text_frame
            tf.margin_left = tf.margin_right = 0
            tf.margin_top = tf.margin_bottom = 0
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = title
            run.font.size = Pt(24)
            run.font.bold = True

        # 6. Place the PNG
        slide.shapes.add_picture(
            str(tmp_png),
            Inches(left_in),
            Inches(top_in),
            Inches(pic_w_in),
            Inches(pic_h_in),
        )

        prs.save(str(pptx_path))
    finally:
        try:
            tmp_png.unlink()
        except OSError:
            pass

    return pptx_path


#Validator — mechanical self-check
# Catches the most common mistakes BEFORE the user sees the output.
# These are deliberately conservative — a "pass" doesn't guarantee
# the diagram is beautiful, only that the obvious failure modes are absent.

SVG_NS = "{http://www.w3.org/2000/svg}"

def validate_svg(svg_path: str | Path) -> list[str]:
    """Return a list of problems. Empty list means the SVG passes basic checks."""
    svg_path = Path(svg_path)
    problems: list[str] = []
    try:
        tree = ET.parse(svg_path)
    except ET.ParseError as e:
        return [f"SVG is not valid XML: {e}"]
    root = tree.getroot()

    # --- check viewBox present ---
    vb = root.attrib.get("viewBox")
    if not vb:
        problems.append("No viewBox attribute on <svg>; cannot determine canvas size.")
        return problems

    vb_parts = vb.replace(",", " ").split()
    if len(vb_parts) != 4:
        problems.append(f"viewBox must have 4 numbers, got: {vb!r}")
        return problems
    _, _, vb_w, vb_h = (float(x) for x in vb_parts)

    # --- check font-family declares a CJK fallback if any Chinese chars are used ---
    all_text = "".join(t.text or "" for t in root.iter(SVG_NS + "text"))
    has_cjk = any("\u4e00" <= ch <= "\u9fff" for ch in all_text)
    font_family_decls = []
    for elem in root.iter():
        ff = elem.attrib.get("font-family")
        if ff:
            font_family_decls.append(ff)
    joined_ff = " ".join(font_family_decls).lower()
    if has_cjk and not any(
        cjk_hint in joined_ff
        for cjk_hint in ["noto sans cjk", "simhei", "noto sans sc", "pingfang", "microsoft yahei"]
    ):
        problems.append(
            "Diagram contains CJK characters but no CJK-capable font is declared in any font-family. "
            "Add e.g. 'Noto Sans CJK SC' or 'SimHei' to the root <svg font-family> attribute."
        )

    # --- check arrowhead markers are defined if 'marker-end' is used anywhere ---
    used_markers = set()
    for elem in root.iter():
        for attr in ("marker-end", "marker-start", "marker-mid"):
            v = elem.attrib.get(attr)
            if v and v.startswith("url(#"):
                used_markers.add(v[5:-1])
    defined_marker_ids = {
        m.attrib.get("id") for m in root.iter(SVG_NS + "marker")
        if m.attrib.get("id")
    }
    missing = used_markers - defined_marker_ids
    if missing:
        problems.append(
            f"Arrow markers referenced but not defined: {sorted(missing)}. "
            f"Add a <marker> definition in <defs>."
        )

    # --- check every <text> stays inside viewBox (light bounds check on x/y) ---
    for t in root.iter(SVG_NS + "text"):
        try:
            x = float(t.attrib.get("x", "0"))
            y = float(t.attrib.get("y", "0"))
        except ValueError:
            continue
        if x < 0 or x > vb_w + 50 or y < 0 or y > vb_h + 50:
            txt_preview = (t.text or "").strip()[:40]
            problems.append(
                f"<text> anchor ({x},{y}) is outside the viewBox ({vb_w}x{vb_h}): '{txt_preview}'"
            )

    # --- check every <rect> stays inside viewBox ---
    for r in root.iter(SVG_NS + "rect"):
        try:
            x = float(r.attrib.get("x", "0"))
            y = float(r.attrib.get("y", "0"))
            w = float(r.attrib.get("width", "0"))
            h = float(r.attrib.get("height", "0"))
        except ValueError:
            continue
        if x < -5 or y < -5 or x + w > vb_w + 5 or y + h > vb_h + 5:
            problems.append(
                f"<rect> at ({x},{y},{w}x{h}) extends outside the viewBox ({vb_w}x{vb_h})."
            )

    # --- check that the SVG uses palette colours, not random ones ---
    # (heuristic: collect all fill/stroke values; warn if there are too many distinct hues)
    fills = set()
    for elem in root.iter():
        for attr in ("fill", "stroke"):
            v = elem.attrib.get(attr, "").strip().lower()
            if v.startswith("#") and len(v) == 7:
                fills.add(v)
    if len(fills) > 18:
        problems.append(
            f"SVG uses {len(fills)} distinct colors. Palette discipline expects ≤ ~18. "
            f"Likely you're mixing palettes — consolidate to one palette from palettes.md."
        )

    return problems



def _main():
    ap = argparse.ArgumentParser(description="Frame-skill render helpers.")
    sub = ap.add_subparsers(dest="cmd", required=True)

    p_val = sub.add_parser("validate", help="Run mechanical self-check on an SVG.")
    p_val.add_argument("svg")

    p_png = sub.add_parser("to_png", help="Render SVG to PNG at given DPI.")
    p_png.add_argument("svg")
    p_png.add_argument("png")
    p_png.add_argument("--dpi", type=int, default=400)

    p_ppt = sub.add_parser("to_pptx", help="Embed SVG into a single-slide PPTX, centered.")
    p_ppt.add_argument("svg")
    p_ppt.add_argument("pptx")
    p_ppt.add_argument("--dpi", type=int, default=400)
    p_ppt.add_argument("--title", default=None)

    p_all = sub.add_parser("to_all", help="Render to both PNG and PPTX next to the SVG.")
    p_all.add_argument("svg")
    p_all.add_argument("--dpi", type=int, default=400)
    p_all.add_argument("--title", default=None)

    args = ap.parse_args()

    if args.cmd == "validate":
        issues = validate_svg(args.svg)
        if not issues:
            print(f"OK: {args.svg} passed all checks.")
            sys.exit(0)
        print(f"FAIL: {len(issues)} issue(s) in {args.svg}:")
        for i, msg in enumerate(issues, 1):
            print(f"  {i}. {msg}")
        sys.exit(1)

    if args.cmd == "to_png":
        out = svg_to_png(args.svg, args.png, dpi=args.dpi)
        print(f"wrote {out}")
        return

    if args.cmd == "to_pptx":
        out = svg_to_pptx(args.svg, args.pptx, dpi=args.dpi, title=args.title)
        print(f"wrote {out}")
        return

    if args.cmd == "to_all":
        base = Path(args.svg).with_suffix("")
        png = svg_to_png(args.svg, f"{base}.png", dpi=args.dpi)
        pptx = svg_to_pptx(args.svg, f"{base}.pptx", dpi=args.dpi, title=args.title)
        print(f"wrote {png}")
        print(f"wrote {pptx}")
        return

if __name__ == "__main__":
    _main()
